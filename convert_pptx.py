#!/usr/bin/env python3
import os
import sys
from pptx import Presentation

# Create output directory
output_dir = "text_version"
os.makedirs(output_dir, exist_ok=True)

def pptx_to_text(pptx_path, txt_path):
    """Extracts text from a .pptx file and saves it as a .txt file."""
    prs = Presentation(pptx_path)
    
    # Extract text from slides
    text = "\n".join(
        shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")
    )

    # Write to output text file
    with open(txt_path, "w", encoding="utf-8") as f:
        f.write(text)

    print(f"Converted: {pptx_path} -> {txt_path}")

def main():
    """Finds all .pptx files in the current directory and converts them."""
    pptx_files = [f for f in os.listdir('.') if f.endswith('.pptx')]

    if not pptx_files:
        print("No .pptx files found in the current directory.")
        sys.exit(1)

    for pptx in pptx_files:
        filename = os.path.splitext(pptx)[0]  # Remove .pptx extension
        txt_path = os.path.join(output_dir, f"{filename}-pptx.txt")
        pptx_to_text(pptx, txt_path)

    print("All PPTX files have been converted!")

if __name__ == "__main__":
    main()

